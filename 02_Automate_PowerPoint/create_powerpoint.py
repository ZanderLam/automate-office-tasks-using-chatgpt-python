import os
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Determine the directory of the script
script_dir = os.path.dirname(os.path.abspath(__file__))

input_folder = os.path.join(script_dir, 'input')
charts_folder = os.path.join(script_dir, 'charts')
ppt_file = 'financial_data.pptx'

# Create the charts folder if it doesn't exist
if not os.path.exists(charts_folder):
    os.mkdir(charts_folder)

# Create a new PowerPoint presentation
prs = Presentation()

# Define a color palette for the charts
colors = sns.color_palette("viridis", as_cmap=True)

# Iterate through all Excel files in the input folder
for excel_file in os.listdir(input_folder):
    if not excel_file.endswith('.xlsx'):
        continue

    # Construct the file path
    file_path = os.path.join(input_folder, excel_file)

    try:
        # Read the financial data from the first worksheet of the Excel file
        df = pd.read_excel(file_path, sheet_name=0, usecols="A:P")

        # Ensure necessary columns exist
        if 'Product' not in df.columns or 'Sales' not in df.columns:
            print(f"Missing 'Product' or 'Sales' column in {excel_file}. Skipping file.")
            continue

        # Drop rows where either 'Product' or 'Sales' is NaN
        df = df[['Product', 'Sales']].dropna()

        # Convert 'Sales' to numeric, forcing errors to NaN
        df['Sales'] = pd.to_numeric(df['Sales'], errors='coerce')
        df = df.dropna()

        # Group the data by the 'Product' column and sum up the 'Sales' column
        grouped = df.groupby('Product').sum()['Sales']

        # Create a chart using the seaborn library
        plt.figure(figsize=(10, 6))
        ax = sns.barplot(x=grouped.index, y=grouped.values, palette='viridis')
        plt.title(f'Sales by Product - {excel_file}')
        plt.xlabel('Product')
        plt.ylabel('Sales')
        plt.xticks(rotation=45)  # Rotate x-axis labels for better readability
        plt.tight_layout()

        # Add sales numbers on top of each bar
        for p in ax.patches:
            ax.annotate(format(p.get_height(), '.2f'), 
                        (p.get_x() + p.get_width() / 2., p.get_height()), 
                        ha='center', va='center', 
                        xytext=(0, 8), 
                        textcoords='offset points')

        # Save the chart to the charts folder
        chart_file = excel_file.replace('.xlsx', '.png')
        chart_path = os.path.join(charts_folder, chart_file)
        plt.savefig(chart_path)
        plt.close()  # Close the plot to avoid overlap

        # Add a slide to the PowerPoint presentation and insert the chart and title
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = excel_file.replace('.xlsx', '')

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(6)
        slide.shapes.add_picture(chart_path, left, top, width=width, height=height)

    except PermissionError as e:
        print(f"Permission error: {e}")
    except Exception as e:
        print(f"An error occurred with file {excel_file}: {e}")

# Save the PowerPoint presentation in the same directory as the script
ppt_path = os.path.join(script_dir, ppt_file)
prs.save(ppt_path)

print(f"PowerPoint presentation saved to {ppt_path}")

# py -m pip install python-pptx
# py -m pip install seaborn
# py -m pip list
# cd "C:\Users\ofschlam\OneDrive - Aibel\Desktop\Github projects\automate-office-tasks-using-chatgpt-python\02_Automate_PowerPoint"
# py create_powerpoint.py